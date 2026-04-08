"""
Build Colorado_Powracle_AlpineFrost.pptx — full 16-slide presentation.

Run: /opt/anaconda3/envs/powracle/bin/python slides/build_deck.py
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ── AlpineFrost palette ──────────────────────────────────────────────────────
ACCENT     = RGBColor(0x7d, 0xc1, 0xe8)   # icy blue
BACKGROUND = RGBColor(0xf0, 0xf7, 0xfb)   # near-white
BODY       = RGBColor(0x2c, 0x42, 0x56)   # dark slate
ALT        = RGBColor(0xdc, 0xea, 0xf2)   # alt row
DARK       = RGBColor(0x1a, 0x2c, 0x3e)   # darker accent for emphasis
MUTED      = RGBColor(0x6b, 0x82, 0x9a)   # muted body

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "slides" / "assets"
OUT = ROOT / "slides" / "Colorado_Powracle_AlpineFrost.pptx"

# ── Helpers ──────────────────────────────────────────────────────────────────
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # Blank

def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill(s, color)
    return s

def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=BODY, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return box

def add_bullets(slide, left, top, width, height, items, *,
                size=16, color=BODY, bullet_color=None, line_spacing=1.15):
    bullet_color = bullet_color or ACCENT
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        if isinstance(item, tuple):
            label, body = item
            r1 = p.add_run(); r1.text = "▸ "; r1.font.size = Pt(size); r1.font.color.rgb = bullet_color; r1.font.bold = True
            r2 = p.add_run(); r2.text = label; r2.font.size = Pt(size); r2.font.color.rgb = DARK; r2.font.bold = True
            r3 = p.add_run(); r3.text = " — " + body; r3.font.size = Pt(size); r3.font.color.rgb = color
        else:
            r1 = p.add_run(); r1.text = "▸ "; r1.font.size = Pt(size); r1.font.color.rgb = bullet_color; r1.font.bold = True
            r2 = p.add_run(); r2.text = item; r2.font.size = Pt(size); r2.font.color.rgb = color
    return box

def page_chrome(slide, title, subtitle=None, eyebrow=None):
    """Standard header bar + title block."""
    # background
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BACKGROUND)
    # accent bar
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), ACCENT)
    # eyebrow (small)
    if eyebrow:
        add_text(slide, Inches(0.6), Inches(0.32), Inches(12), Inches(0.3),
                 eyebrow, size=11, bold=True, color=ACCENT)
    # title
    add_text(slide, Inches(0.6), Inches(0.55), Inches(12.5), Inches(0.7),
             title, size=30, bold=True, color=DARK)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.15), Inches(12.5), Inches(0.4),
                 subtitle, size=15, color=MUTED)

def add_note(slide, text, *, top=None):
    """Sticky-note style annotation for owner handoffs."""
    top = top or Inches(6.6)
    n = add_rect(slide, Inches(0.6), top, Inches(12.1), Inches(0.6), ALT)
    n.line.color.rgb = ACCENT
    n.line.width = Pt(1)
    add_text(slide, Inches(0.75), top + Inches(0.08), Inches(11.9), Inches(0.5),
             text, size=11, color=BODY, bold=False)

def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

# ── Build ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W, H = prs.slide_width, prs.slide_height

# 1. TITLE -------------------------------------------------------------------
s = add_blank(prs)
add_rect(s, 0, 0, W, H, BACKGROUND)
add_rect(s, 0, 0, W, Inches(4.2), ACCENT)
add_text(s, Inches(0.8), Inches(1.4), Inches(12), Inches(1.2),
         "Colorado Powracle", size=58, bold=True, color=RGBColor(0xff,0xff,0xff))
add_text(s, Inches(0.8), Inches(2.5), Inches(12), Inches(0.8),
         "A LangChain ski-conditions agent grounded in real Colorado data",
         size=22, color=RGBColor(0xff,0xff,0xff))
add_text(s, Inches(0.8), Inches(4.6), Inches(12), Inches(0.5),
         "Big Data Architecture · Spring 2026", size=16, color=MUTED, bold=True)
add_text(s, Inches(0.8), Inches(5.1), Inches(12), Inches(0.5),
         "[ team names — fill in ]", size=16, color=BODY)
add_text(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
         "github.com/Ttopiac/colorado-powracle", size=12, color=MUTED)
add_speaker_notes(s, "Title slide. Introduce team and one-line pitch.")

# 2. THE PROBLEM -------------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "The Problem", eyebrow="WHY POWRACLE")
add_bullets(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(2.5), [
    "Colorado skiers juggle snow reports, I-70 traffic, weekend forecasts, and pass eligibility across 19 resorts.",
    "No single tool answers the real question:  \u201cwhere should I ski this weekend?\u201d",
    "We wanted answers grounded in real data — not LLM guesses.",
], size=18)
# hook
add_rect(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(0.65), ACCENT)
add_text(s, Inches(0.75), Inches(4.46), Inches(12), Inches(0.55),
         "So… where does the grounded data come from?",
         size=20, bold=True, color=RGBColor(0xff,0xff,0xff))
# answer
add_bullets(s, Inches(0.6), Inches(5.25), Inches(12.1), Inches(2), [
    ("Live APIs", "SNOTEL · COtrip · Open-Meteo · SerpAPI"),
    ("10 yr Historical Snow", "SNOTEL daily CSVs → Parquet → DuckDB (40,780 rows)"),
    ("10 yr Historical Traffic", "CDOT hourly volumes → DuckDB traffic_patterns view"),
], size=16)
add_speaker_notes(s, "Set up the question the rest of the deck answers. Emphasize 'grounded data'.")

# 3. WHAT WE BUILT + SIMPLE OVERVIEW -----------------------------------------
s = add_blank(prs)
page_chrome(s, "What We Built", "ReAct agent · 6 tools · Streamlit UI + FastAPI · grounded data", eyebrow="OVERVIEW")
img = ASSETS / "simple_overview.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(1.4), Inches(1.7), width=Inches(10.5))
add_text(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5),
         "Streamlit UI and FastAPI both call the same chat_service.run_chat_turn — one agent, two front doors.",
         size=13, color=MUTED, bold=True)
add_speaker_notes(s, "Simple overview. Mention dual entry points share one chat service. Foreshadow deterministic fast path.")

# 4. THE AGENT — REACT LOOP --------------------------------------------------
s = add_blank(prs)
page_chrome(s, "The Agent — ReAct Loop", "Claude 3 Haiku via OpenRouter · zero-shot-react-description", eyebrow="AGENT")
# left: loop diagram boxes
def loop_box(left, top, label, color):
    box = add_rect(s, left, top, Inches(2.6), Inches(0.85), color)
    add_text(s, left, top + Inches(0.18), Inches(2.6), Inches(0.5),
             label, size=15, bold=True, color=RGBColor(0xff,0xff,0xff), align=PP_ALIGN.CENTER)
positions = [
    (Inches(0.7), Inches(2.0), "Question", ACCENT),
    (Inches(0.7), Inches(3.1), "Thought", DARK),
    (Inches(0.7), Inches(4.2), "Action  (tool call)", DARK),
    (Inches(0.7), Inches(5.3), "Observation", DARK),
    (Inches(0.7), Inches(6.4), "Final Answer", ACCENT),
]
for l,t,lbl,c in positions: loop_box(l,t,lbl,c)
# arrow text
add_text(s, Inches(3.5), Inches(3.2), Inches(2), Inches(0.4), "↓", size=22, bold=True, color=ACCENT)
add_text(s, Inches(3.5), Inches(4.3), Inches(2), Inches(0.4), "↓", size=22, bold=True, color=ACCENT)
add_text(s, Inches(3.5), Inches(5.4), Inches(2), Inches(0.4), "↺ repeat", size=14, bold=True, color=ACCENT)

# right: explanation
add_bullets(s, Inches(4.6), Inches(2.0), Inches(8.2), Inches(5), [
    ("LLM reasons step-by-step", "picks the right tool per sub-question"),
    ("Multi-tool composition", "snow + traffic + forecast in one turn"),
    ("Stop condition", "agent emits Final Answer when it has enough evidence"),
    ("Mandatory output contract", "every reply ends with [RANKING: …] line parsed by the UI"),
    ("Why ReAct over function-call only?", "the Thought trace is debuggable and shows up in our verbose logs"),
], size=16)
add_speaker_notes(s, "Set up ReAct mental model before showing the prompt.")

# 5. PROMPT ARCHITECTURE -----------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Prompt Architecture", "agent/prompts.py — the agent's instruction manual", eyebrow="HOW WE STEER THE AGENT")

# Left half: stylised code-window screenshot rendered as a shape
code_left = Inches(0.6); code_top = Inches(1.7)
code_w = Inches(6.0); code_h = Inches(5.2)
window = add_rect(s, code_left, code_top, code_w, code_h, RGBColor(0x1a,0x2c,0x3e))
# title bar
add_rect(s, code_left, code_top, code_w, Inches(0.35), RGBColor(0x12,0x22,0x32))
add_text(s, code_left + Inches(0.15), code_top + Inches(0.05), Inches(5.7), Inches(0.3),
         "agent/prompts.py", size=11, color=RGBColor(0x7d,0xc1,0xe8), font="Menlo", bold=True)

code_lines = [
    ('SYSTEM_PROMPT = """', RGBColor(0xc6,0xa5,0xff)),
    ('You are the Colorado Powder Oracle — concise,', RGBColor(0xe6,0xee,0xf6)),
    ('knowledgeable assistant for Colorado skiers.', RGBColor(0xe6,0xee,0xf6)),
    ('', BODY),
    ('CONVERSATION MEMORY:', RGBColor(0x7d,0xc1,0xe8)),
    ('  follow-ups → previous topic; new topic →', RGBColor(0xe6,0xee,0xf6)),
    ('  treat as fresh question (topic-bleed fix)', RGBColor(0xe6,0xee,0xf6)),
    ('', BODY),
    ('LIVE DATA ALREADY IN CONTEXT:', RGBColor(0x7d,0xc1,0xe8)),
    ('  read snapshot before calling tools', RGBColor(0xe6,0xee,0xf6)),
    ('', BODY),
    ('TOOLS AND WHEN TO USE THEM:', RGBColor(0x7d,0xc1,0xe8)),
    ('  "historically", "average", month names →', RGBColor(0xe6,0xee,0xf6)),
    ('  get_snowpack_history (never current)', RGBColor(0xe6,0xee,0xf6)),
    ('  weekend / forecast → get_snow_forecast', RGBColor(0xe6,0xee,0xf6)),
    ('  road / chains → get_live_traffic', RGBColor(0xe6,0xee,0xf6)),
    ('', BODY),
    ('RESORT KNOWLEDGE:  IKON · EPIC · INDY', RGBColor(0x7d,0xc1,0xe8)),
    ('SNOWPACK SCIENCE:   SWE · powder thresholds', RGBColor(0x7d,0xc1,0xe8)),
    ('TRAFFIC KNOWLEDGE:  I-70 peaks · Eisenhower', RGBColor(0x7d,0xc1,0xe8)),
    ('', BODY),
    ('MANDATORY RANKING LINE:', RGBColor(0x7d,0xc1,0xe8)),
    ('  [RANKING: Resort1, Resort2, ...]', RGBColor(0xe6,0xee,0xf6)),
    ('"""', RGBColor(0xc6,0xa5,0xff)),
]
code_box = s.shapes.add_textbox(code_left + Inches(0.2), code_top + Inches(0.45),
                                 code_w - Inches(0.4), code_h - Inches(0.55))
tf = code_box.text_frame; tf.word_wrap = True
for i,(line,clr) in enumerate(code_lines):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.line_spacing = 1.05
    r = p.add_run(); r.text = line or " "
    r.font.size = Pt(11); r.font.name = "Menlo"; r.font.color.rgb = clr

# Right half: section breakdown
add_text(s, Inches(6.95), Inches(1.7), Inches(6.2), Inches(0.4),
         "What's actually inside (≈80 lines):", size=14, bold=True, color=DARK)
add_bullets(s, Inches(6.95), Inches(2.1), Inches(6.2), Inches(5), [
    ("1. Conversation memory rules", "follow-ups vs. topic switches"),
    ("2. Live data already in context", "read snapshot first, don't loop tools"),
    ("3. Tool routing rules", "keyword → exact tool"),
    ("4. Resort knowledge", "19 resorts grouped by pass"),
    ("5. Snowpack science + traffic", "SWE, I-70 peak hours, chokepoints"),
    ("6. Output contract", "[RANKING: …] line parsed by the UI"),
], size=14)
add_text(s, Inches(6.95), Inches(6.55), Inches(6.2), Inches(0.6),
         "The prompt is a routing table, a domain knowledge base, and an output contract — all in one.",
         size=12, color=ACCENT, bold=True)
add_speaker_notes(s, "Walk the right column. Mention this is also where the topic-bleed fix lives — full bug story on the lessons slide.")

# 6. THE 6 TOOLS -------------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "The 6 Tools", "Each tool is a typed function the agent can call by name", eyebrow="TOOLS")

# Table-ish grid
headers = ["Tool", "Source", "Latency"]
rows = [
    ("get_current_snowpack",     "SNOTEL REST API (live)",          "~2 s"),
    ("get_snowpack_history",     "DuckDB · 10 yr daily SNOTEL",     "~50 ms"),
    ("get_live_traffic",         "COtrip API (live)",                "~1 s"),
    ("get_best_departure_time",  "DuckDB · traffic_patterns view",   "~50 ms"),
    ("get_snow_forecast",        "Open-Meteo (HRRR → best_match)",   "~1 s"),
    ("web_search",               "SerpAPI (fallback for anything)", "~2 s"),
]
left = Inches(0.6); top = Inches(1.7)
col_w = [Inches(4.2), Inches(5.6), Inches(2.7)]
row_h = Inches(0.55)
# header
x = left
for i,h in enumerate(headers):
    add_rect(s, x, top, col_w[i], row_h, ACCENT)
    add_text(s, x + Inches(0.15), top + Inches(0.12), col_w[i] - Inches(0.2), Inches(0.4),
             h, size=14, bold=True, color=RGBColor(0xff,0xff,0xff))
    x += col_w[i]
# rows
for ri, row in enumerate(rows):
    y = top + row_h + ri * row_h
    bg = ALT if ri % 2 == 0 else BACKGROUND
    x = left
    for ci, val in enumerate(row):
        add_rect(s, x, y, col_w[ci], row_h, bg)
        add_text(s, x + Inches(0.15), y + Inches(0.12), col_w[ci] - Inches(0.2), Inches(0.4),
                 val, size=13, color=BODY, bold=(ci==0),
                 font=("Menlo" if ci==0 else "Calibri"))
        x += col_w[ci]

add_text(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(0.5),
         "Dual data path:", size=14, bold=True, color=DARK)
add_bullets(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.3), [
    "UI bulk-fetches SNOTEL + forecasts on page load (cached 30 min / 3 hr)",
    "Agent calls the same tools per conversation — no double work, fresh data when needed",
], size=13)
add_speaker_notes(s, "Don't read the whole table aloud. Highlight live vs DuckDB split and the web_search safety net.")

# 7. LIVE WALKTHROUGH — DRONA -------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Live Walkthrough — Real ReAct Trace", "Captured from a real agent run, not a mock-up", eyebrow="DEMO")

add_text(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(0.4),
         "Suggested prompt (forces multiple tools):", size=14, bold=True, color=DARK)
quote = add_rect(s, Inches(0.6), Inches(2.1), Inches(12.1), Inches(0.7), ALT)
quote.line.color.rgb = ACCENT
add_text(s, Inches(0.8), Inches(2.22), Inches(11.8), Inches(0.5),
         '"I want powder this Saturday but avoid I-70 — when should I leave Denver?"',
         size=16, color=DARK, bold=True)

add_text(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(0.4),
         "Expected trace shape:", size=14, bold=True, color=DARK)
add_bullets(s, Inches(0.6), Inches(3.5), Inches(12.1), Inches(2.5), [
    ("Thought 1", "user wants powder + traffic-aware departure → call snow forecast"),
    ("Action 1",  "get_snow_forecast(\"Steamboat Springs\") · get_snow_forecast(\"Winter Park\")"),
    ("Thought 2", "now resolve traffic → check best departure for non-I-70 corridor"),
    ("Action 2",  "get_best_departure_time(\"US-40 Saturday\")"),
    ("Final",     "ranked recommendation + departure window + [RANKING: …]"),
], size=13)

add_note(s, "DRONA — please rewrite this slide with a real captured trace. "
            "Run the agent with verbose=True against the prompt above (or your own). "
            "Show Question → Thought / Action / Observation cycles → Final Answer.")
add_speaker_notes(s, "Drona owns this slide. Goal: a single multi-tool real trace, not invented.")

# 8. DETERMINISTIC FAST PATH -------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Deterministic Fast Path", "agent/deterministic_answers.py — opt-in shortcut", eyebrow="OPTIMIZATION")

add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.5), [
    ("What it covers", "exactly two metrics — most fresh snow (72 h) and deepest base — nothing else"),
    ("Why", "for these questions, reasoning adds no value; SNOTEL data is already structured"),
    ("Toggle", "checkbox in the UI sidebar, flag in the FastAPI request body"),
], size=15)

# Decision flow
add_text(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.4),
         "Decision logic:", size=14, bold=True, color=DARK)

def flow_box(left, top, w, label, color, txt_color=None):
    add_rect(s, left, top, w, Inches(0.85), color)
    add_text(s, left, top + Inches(0.2), w, Inches(0.5),
             label, size=12, bold=True,
             color=txt_color or RGBColor(0xff,0xff,0xff), align=PP_ALIGN.CENTER)

flow_box(Inches(0.6),  Inches(4.15), Inches(2.4), "Question in", ACCENT)
flow_box(Inches(3.2),  Inches(4.15), Inches(2.6), "Blocker words?", DARK)
flow_box(Inches(6.0),  Inches(4.15), Inches(2.6), "Trigger phrase?", DARK)
flow_box(Inches(8.8),  Inches(4.15), Inches(3.9), "Sort 19 resorts → top 1 + next 2", ACCENT)
# arrows
for x in (Inches(3.0), Inches(5.8), Inches(8.6)):
    add_text(s, x, Inches(4.4), Inches(0.3), Inches(0.4), "→", size=18, bold=True, color=ACCENT)

# fall through note
add_text(s, Inches(3.2), Inches(5.15), Inches(5.4), Inches(0.4),
         "yes → return None → fall back to LLM", size=12, color=MUTED, bold=True)

add_text(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(0.4),
         "Blocker words:", size=13, bold=True, color=DARK)
add_text(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.4),
         "historical · average · january…december · forecast · weekend · recommend · compare · consistent",
         size=12, color=BODY, font="Menlo")

add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.4),
         "Trigger phrases:", size=13, bold=True, color=DARK)
add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4),
         '"fresh snow" · "new snow" · "snowiest right now" · "deepest base" · "biggest base"',
         size=12, color=BODY, font="Menlo")
add_speaker_notes(s, "Fast path is intentionally narrow. Anything ambiguous falls through to the agent untouched.")

# 9. EVALUATION — DRONA -------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Evaluation Harness", "Built around the deterministic fast path", eyebrow="MEASURING THE AGENT")
img = ASSETS / "evaluation_harness.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.6), Inches(1.7), width=Inches(8.0))
add_bullets(s, Inches(8.9), Inches(1.8), Inches(4.3), Inches(5), [
    "30 benchmark prompts",
    "10 factual",
    "10 recommendation",
    "10 explanatory",
    "Compares deterministic fast-path vs full agent",
    "Baselines committed in eval/baselines/",
], size=13)
add_note(s, "DRONA — please rewrite. Add real numbers from the latest run (factual accuracy, "
            "deterministic vs agent, latency). Code: eval/run_agent_eval.py · score_outputs.py · plot_results.py.")
add_speaker_notes(s, "Drona owns. Show one bar chart from eval/figures/ if possible.")

# 10. STREAMLIT UI ------------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Streamlit UI", "What the user actually sees", eyebrow="FRONT END")
shot = ASSETS / "ui_map.png"
if shot.exists():
    s.shapes.add_picture(str(shot), Inches(0.6), Inches(1.7), width=Inches(7.5))

add_bullets(s, Inches(8.3), Inches(1.7), Inches(5.0), Inches(5.0), [
    ("Map", "Plotly Scattermapbox · ESRI topo · pass-colored markers"),
    ("Today's Leaders", "fresh snow · base · closest powder banner"),
    ("Quick filters", "6\"+ powder · 50\"+ base · <100 mi · 4\"+ weekend"),
    ("Sort modes", "Fresh · Base · Distance · AI Pick"),
    ("Smart Trip Planner", "multi-day itinerary, distances injected"),
    ("Snowfall effect", "CSS animation (proposed: default on)"),
    ("Theme", "Mountain Stone — slate + accent blue"),
], size=12)
add_note(s, "JOSEF — please review. Suggestion: drop the snowfall toggle and just default the snowfall "
            "effect on, if the CSS animation isn't expensive on lower-end machines. Add anything else "
            "you'd like to highlight on the UI.")
add_speaker_notes(s, "Will capture real screenshots before final.")

# 11. FASTAPI -----------------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "FastAPI Service", "Programmatic access — same agent, same chat service", eyebrow="API LAYER")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4), [
    ("POST /chat", "ChatRequest → ChatResponse with answer, ranking, raw_response"),
    ("GET /health", "liveness probe"),
    ("Same chat_service.run_chat_turn as the Streamlit UI", "no logic duplication"),
    ("TTL caches", "30 min for live conditions, 3 hr for forecasts (matches UI)"),
    ("Lazy agent init", "first /chat builds the agent — /health stays cheap"),
    ("Why it matters", "future mobile client, third-party integrations, eval harness uses it too"),
], size=15)
add_note(s, "JOSEF — please rewrite. Code: api.py. Feel free to add request/response examples or curl snippets.")
add_speaker_notes(s, "Josef owns.")

# 12. USER ACCOUNTS -----------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "User Accounts & Pass ROI", "PostgreSQL-backed personalisation", eyebrow="ACCOUNTS")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.5), [
    ("Login / register", "auth/auth_manager.py — sessions in PostgreSQL"),
    ("Profile", "home city, ski ability, terrain preference, ski pass management"),
    ("Trips", "trip history with per-day check-ins and star ratings"),
    ("Stats", "season ROI dashboard — days skied, pass cost, ticket value, break-even progress"),
    ("ROI calculation", "auth/roi_calculator.py · day_ticket_price tracked per pass"),
    ("Guest mode", "if Postgres is unavailable, login UI hides and core features still work"),
], size=15)
add_note(s, "JOSEF — please rewrite. Code: auth/, models/user.py, db/postgres.py, db/run_migrations.py.")
add_speaker_notes(s, "Josef owns.")

# 13. ENGINEERING LESSONS / REAL BUGS ----------------------------------------
s = add_blank(prs)
page_chrome(s, "Engineering Lessons — Real Bugs We Hit", "Every fix is in the git history", eyebrow="WHAT BROKE & WHY")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.6), [
    ("Topic-bleed in conversation memory  (commit 6e816f0)",
     "Agent kept answering the previous question after a topic switch. Asked about resorts, then asked about restaurants → still got resorts. Fixed with an explicit \u201ctreat NEW topics as fresh\u201d rule in SYSTEM_PROMPT."),
    ("Historical questions hit the wrong tool  (commit 9f7ab6a)",
     "Agent called get_current_snowpack when asked \u201cmost snow in January\u201d. Fixed with explicit keyword routing in the prompt — \u201chistorically\u201d, \u201con average\u201d, month names → get_snowpack_history."),
    ("SNOTEL API surprises  (CLAUDE.md notes)",
     "Network code is SNTL not SNOTEL. /data response nests elements under a per-station data key. Both broke our first ingestion attempt."),
    ("AI Pick poisoned by assistant replies  (commit af1a6d1)",
     "We were passing assistant responses as ranking context, biasing later rankings. Fixed by passing only user messages."),
    ("AI Pick silent failure  (commit ffa2464)",
     "Used deprecated .predict(); switched to .invoke(), surfaced errors instead of falling back silently."),
    ("Sequential SNOTEL fetch killed page load  (commit 9db93df)",
     "Replaced 19 sequential calls with batched + progressive loading."),
], size=13)
add_speaker_notes(s, "Pick the top 3 to actually narrate. Rest are backup if asked. All grounded in commits.")

# 14. AI-ASSISTED COLLABORATION ----------------------------------------------
s = add_blank(prs)
page_chrome(s, "AI-Assisted Collaboration", "How three humans + AI assistants shipped this without breaking each other", eyebrow="WORKFLOW")

# left: stack of three cards
def card(left, top, w, h, title, body):
    add_rect(s, left, top, w, h, ALT).line.color.rgb = ACCENT
    add_text(s, left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.4),
             title, size=13, bold=True, color=DARK)
    add_text(s, left + Inches(0.15), top + Inches(0.45), w - Inches(0.3), h - Inches(0.5),
             body, size=11, color=BODY)

cw = Inches(6.1); cl = Inches(0.6)
card(cl, Inches(1.7), cw, Inches(1.55),
     "CLAUDE.md  ·  distributed context",
     "Root file = project overview, run commands, the 6 tools, key decisions.\n"
     "Subdir CLAUDE.md in agent/, tools/, ingestion/, db/ = local gotchas the AI picks up automatically.")
card(cl, Inches(3.35), cw, Inches(1.55),
     "AGENTS.md  ·  cross-tool sync",
     "Same content as root CLAUDE.md but in the convention used by Cursor and other AI tools.\n"
     "Kept in sync via a checkbox in the PR template — every assistant sees the same rules.")
card(cl, Inches(5.0), cw, Inches(1.95),
     ".github/pull_request_template.md  +  /review-pr skill",
     "Template auto-loads on every PR: streamlit run cleanly · 2 canonical questions tested · "
     "no secrets staged · docs updated in the same PR · tools registered + prompt updated.\n"
     "Custom Claude Code skill .claude/skills/review-pr/SKILL.md walks every checkbox against the "
     "actual diff and runtime (not just whether the box is ticked) and posts an audit comment.")

# right: flow diagram
fl = Inches(7.0); fw = Inches(6.0)
def flow(top, label, color):
    add_rect(s, fl, top, fw, Inches(0.6), color)
    add_text(s, fl, top + Inches(0.13), fw, Inches(0.4),
             label, size=13, bold=True, color=RGBColor(0xff,0xff,0xff), align=PP_ALIGN.CENTER)

flow(Inches(1.8), "Developer + AI assistant", ACCENT)
flow(Inches(2.6), "reads CLAUDE.md / AGENTS.md", DARK)
flow(Inches(3.4), "writes code", DARK)
flow(Inches(4.2), "opens PR  →  checklist auto-loads", DARK)
flow(Inches(5.0), "/review-pr skill audits diff + runtime", DARK)
flow(Inches(5.8), "merge to main", ACCENT)
for top in (Inches(2.4), Inches(3.2), Inches(4.0), Inches(4.8), Inches(5.6)):
    add_text(s, fl + Inches(2.85), top, Inches(0.4), Inches(0.3), "↓", size=14, bold=True, color=ACCENT)

add_text(s, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.4),
         "The repo is its own onboarding doc — every bug we hit became a rule the next contributor (human or AI) inherits.",
         size=12, bold=True, color=ACCENT)
add_speaker_notes(s, "Tie this back to slide 13: bugs become rules in CLAUDE.md, prevented from recurring.")

# 15. FUTURE WORK -------------------------------------------------------------
s = add_blank(prs)
page_chrome(s, "Future Work", "Where we'd take Powracle next", eyebrow="ROADMAP")
add_bullets(s, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5), [
    ("Phase 3 — Location & routing", "OpenRouteService for real driving distances and ETAs"),
    ("Mobile client", "thin React Native app over the existing FastAPI"),
    ("Larger eval set", "expand beyond 30 prompts; nightly regression against eval/baselines/"),
    ("More deterministic shortcuts", "extend the fast path to other narrow factual classes"),
    ("Personalised ranking", "feed user pass + ski-day history into the agent prompt"),
], size=16)
add_speaker_notes(s, "Optional slide — drop if running long.")

# 16. THANK YOU ---------------------------------------------------------------
s = add_blank(prs)
add_rect(s, 0, 0, W, H, BACKGROUND)
add_rect(s, 0, Inches(2.0), W, Inches(3.5), ACCENT)
add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(1.2),
         "Thank You", size=64, bold=True, color=RGBColor(0xff,0xff,0xff))
add_text(s, Inches(0.8), Inches(3.7), Inches(12), Inches(0.6),
         "Questions?", size=22, color=RGBColor(0xff,0xff,0xff))
add_text(s, Inches(0.8), Inches(5.8), Inches(12), Inches(0.5),
         "[ team names — fill in ]", size=16, color=BODY, bold=True)
add_text(s, Inches(0.8), Inches(6.25), Inches(12), Inches(0.4),
         "github.com/Ttopiac/colorado-powracle", size=13, color=MUTED)
add_text(s, Inches(0.8), Inches(6.65), Inches(12), Inches(0.4),
         "Palette: AlpineFrost", size=11, color=MUTED)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"wrote {OUT}  ({OUT.stat().st_size:,} bytes)  {len(prs.slides)} slides")
