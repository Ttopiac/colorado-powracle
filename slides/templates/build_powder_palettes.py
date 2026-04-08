#!/usr/bin/env python3
"""Build 3-slide palette samples inspired by colorado_powder_oracle.

Replaces the beige/khaki Moderna deck colors with cold, snowy light-blue
palettes. The colorado_powder_oracle Streamlit theme uses:
    primaryColor              = #63b3ed
    backgroundColor           = #383f4a   (Mountain Stone dark)
    secondaryBackgroundColor  = #424e5c
    textColor                 = #dceeff

Each palette produces its own .pptx with 3 pages: Title, Content, Thank You.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

OUT_DIR = "/Users/chli4608/Repositories/resumes-and-cover-letters/interview_prep/palette_samples"


def hx(s):
    s = s.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ---- Palette definitions: cold / snowy / powder-day ----
PALETTES = [
    {
        "name": "PowderDay_Light",
        "tag":  "Powder Day  —  light bluebird",
        "accent":   "#63b3ed",  # CPO primary
        "dark":     "#1e3a5f",
        "bg":       "#f4f9ff",
        "alt":      "#e1efff",
        "title_tx": "#ffffff",
        "subtitle": "#dceeff",
    },
    {
        "name": "AlpineFrost",
        "tag":  "Alpine Frost  —  pale icy blue",
        "accent":   "#7dc1e8",
        "dark":     "#2c4256",
        "bg":       "#f0f7fb",
        "alt":      "#dceaf2",
        "title_tx": "#ffffff",
        "subtitle": "#eaf4ff",
    },
    {
        "name": "MountainStone_Dark",
        "tag":  "Mountain Stone  —  CPO native dark",
        "accent":   "#63b3ed",
        "dark":     "#dceeff",   # text on dark bg
        "bg":       "#383f4a",
        "alt":      "#424e5c",
        "title_tx": "#0f1620",
        "subtitle": "#dceeff",
    },
    {
        "name": "Bluebird_Crisp",
        "tag":  "Bluebird Crisp  —  saturated cold blue",
        "accent":   "#4a90d9",
        "dark":     "#0f2a4a",
        "bg":       "#eaf4ff",
        "alt":      "#d6e8f8",
        "title_tx": "#ffffff",
        "subtitle": "#dceeff",
    },
]


def build_deck(p):
    ACCENT   = hx(p["accent"])
    DARK     = hx(p["dark"])
    BG       = hx(p["bg"])
    ALT      = hx(p["alt"])
    TITLE_TX = hx(p["title_tx"])
    SUBTITLE = hx(p["subtitle"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = BG

    def add_title_bar(slide, text, subtitle=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                       prs.slide_width, Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.2)
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.text = text
        para.font.size = Pt(32); para.font.bold = True
        para.font.color.rgb = TITLE_TX
        para.alignment = PP_ALIGN.LEFT
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(18)
            p2.font.color.rgb = SUBTITLE

    def add_bullets(slide, left, top, w, h, bullets, size=18):
        tb = slide.shapes.add_textbox(left, top, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = b
            para.font.size = Pt(size)
            para.font.color.rgb = DARK
            para.space_after = Pt(6)

    # ---------- SLIDE 1: Title ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                   prs.slide_width, Inches(4.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.8); tf.margin_top = Inches(1.0)
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = "Moderna Co-op Interview Prep"
    para.font.size = Pt(40); para.font.bold = True
    para.font.color.rgb = TITLE_TX
    p2 = tf.add_paragraph()
    p2.text = f"Palette sample: {p['tag']}"
    p2.font.size = Pt(22); p2.font.color.rgb = SUBTITLE
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    p3.text = "Chi-Hui Lin  |  April 2026"
    p3.font.size = Pt(20); p3.font.color.rgb = SUBTITLE
    p3.space_before = Pt(20)

    info = slide.shapes.add_textbox(Inches(0.8), Inches(5.0),
                                    Inches(11.5), Inches(2))
    add_bullets(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2), [
        f"Accent  {p['accent']}      Background  {p['bg']}",
        f"Body text  {p['dark']}      Alt row  {p['alt']}",
        "Cold, snowy bluebird-day palette inspired by colorado_powder_oracle",
    ], size=18)

    # ---------- SLIDE 2: Content ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title_bar(slide, "Spatial Omics Fusion",
                  "GAT + Cross-Attention on DLPFC brain tissue")
    add_bullets(slide, Inches(0.6), Inches(1.7), Inches(12.2), Inches(5), [
        "• MLP encoder: 3,000 HVGs → 128-d expression embedding per spot",
        "• Graph Attention Net on k=6 KNN spatial graph → 128-d spatial embedding",
        "• Cross-attention fusion: each spot attends to most informative neighbors",
        "• Result: ARI 0.926 on DLPFC benchmark, 1.1M params, <30 sec/slice",
        "• Ablation — removing spatial context drops ARI 0.92 → 0.37",
        "• scGPT (33M-cell pretrained, frozen) underperforms task-specific MLP",
    ], size=20)

    # Sample alt-row band so the user can see the secondary color
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(6.3),
                                  Inches(12.2), Inches(0.7))
    band.fill.solid(); band.fill.fore_color.rgb = ALT
    band.line.fill.background()
    tb = band.text_frame
    tb.margin_left = Inches(0.3); tb.margin_top = Inches(0.1)
    bp = tb.paragraphs[0]
    bp.text = "  Secondary surface — alt table rows, callouts, sidebars"
    bp.font.size = Pt(14); bp.font.color.rgb = DARK

    # ---------- SLIDE 3: Thank You ----------
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.0),
                                   prs.slide_width, Inches(3.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.8)
    para = tf.paragraphs[0]
    para.text = "Thank You"
    para.font.size = Pt(44); para.font.bold = True
    para.font.color.rgb = TITLE_TX
    para.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = ""; p2.space_before = Pt(20)
    p3 = tf.add_paragraph()
    p3.text = "Chi-Hui Lin  |  chli4608@colorado.edu"
    p3.font.size = Pt(22); p3.font.color.rgb = SUBTITLE
    p3.alignment = PP_ALIGN.CENTER
    p4 = tf.add_paragraph()
    p4.text = f"Palette: {p['name']}"
    p4.font.size = Pt(16); p4.font.color.rgb = SUBTITLE
    p4.alignment = PP_ALIGN.CENTER

    out = f"{OUT_DIR}/Palette_{p['name']}.pptx"
    prs.save(out)
    print(f"Saved {out}")


if __name__ == "__main__":
    for pal in PALETTES:
        build_deck(pal)
